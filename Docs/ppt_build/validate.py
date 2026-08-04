# -*- coding: utf-8 -*-
"""PowerPoint 없이 덱을 검사한다.
  1) 도형 크기가 음수/0 인가 (→ 파일이 아예 안 열린다)
  2) 패널 밖으로 나갔는가
  3) 글자가 상자를 넘치는가 (렌더 없이 추정)
"""
import sys, io, math, unicodedata
from pptx import Presentation
from pptx.util import Emu

PANEL_TOP, PANEL_BOT = 0.74, 14.27


def char_w(ch, size_pt):
    """대략적인 글자 폭(pt). 한글·한자·가나는 1em, 그 외는 0.52em."""
    if ch == " ":
        return size_pt * 0.30
    ea = unicodedata.east_asian_width(ch)
    return size_pt * (1.0 if ea in ("W", "F") else 0.52)


def est_height(tf, box_w_in):
    """텍스트프레임이 실제로 차지할 높이(inch) 추정."""
    usable_pt = max(1.0, (box_w_in - (Emu(tf.margin_left).inches + Emu(tf.margin_right).inches)) * 72.0)
    total = 0.0
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            continue
        size = max((r.font.size.pt if r.font.size else 18) for r in runs)
        w = sum(char_w(c, (r.font.size.pt if r.font.size else 18)) for r in runs for c in r.text)
        lines = max(1, math.ceil(w / usable_pt))
        ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.2
        sb = p.space_before.pt if p.space_before else 0
        sa = p.space_after.pt if p.space_after else 0
        total += (lines * size * 1.22 * ls + sb + sa) / 72.0
    return total


def main(path):
    prs = Presentation(path)
    fatal = warn = 0
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if sh.width is None or sh.height is None:
                print(f"[치명] 슬라이드{i:2} {sh.name}: 크기 없음"); fatal += 1; continue
            w, h = Emu(sh.width).inches, Emu(sh.height).inches
            t = Emu(sh.top).inches
            is_conn = "Connector" in sh.name
            if w < 0 or h < 0 or (not is_conn and (w == 0 or h == 0)):
                print(f"[치명] 슬라이드{i:2} [{sh.name}] w={w:.2f} h={h:.2f} top={t:.2f}")
                fatal += 1
                continue
            if sh.rotation == 0 and not is_conn:
                if t + h > PANEL_BOT + 0.01 or t < PANEL_TOP - 0.01:
                    print(f"[경고] 슬라이드{i:2} [{sh.name}] 패널 이탈 top={t:.2f} bottom={t+h:.2f}")
                    warn += 1
            if sh.has_text_frame and sh.text_frame.text.strip():
                need = est_height(sh.text_frame, w)
                if need > h * 1.06 + 0.02:
                    txt = sh.text_frame.text.replace("\n", " / ")[:44]
                    print(f"[넘침] 슬라이드{i:2} [{sh.name}] 필요 {need:.2f} > 상자 {h:.2f}  «{txt}»")
                    warn += 1
    print("=" * 60)
    print(f"치명 {fatal}건 · 경고 {warn}건")
    return fatal


if __name__ == "__main__":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
    sys.exit(1 if main(sys.argv[1]) else 0)
