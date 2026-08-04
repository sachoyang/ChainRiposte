# -*- coding: utf-8 -*-
"""
ChainRiposte 포트폴리오 PPT 생성기 (17장)

- _theme_base.pptx(테마 원본)를 복사해 슬라이드만 전부 지우고 새로 짓는다.
  → 임베드된 Pretendard 폰트 · 슬라이드 마스터 · 테마가 그대로 살아 있다.
- 내용은 Docs/ppt대본.md 대본을 그대로 따른다(문장을 지어내지 않는다).
- 세로 배치는 V 배분기가 잡는다. 크기가 음수가 되면 그 자리에서 빌드가 죽는다
  (전에 음수 크기 도형 때문에 PowerPoint가 파일을 아예 못 열었다).
"""
import re, shutil, sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))

# 테마·슬라이드 마스터·임베드 폰트(Pretendard)를 물려받는 원본. 26.67 x 15 inch 캔버스다.
# ⚠ 완성본(Docs/ChainRiposte.pptx)을 여기에 쓰지 말 것 — 구글 슬라이드를 한 번 거치면
#    캔버스가 10 x 5.625 로 줄고 도형 이름이 바뀌어, 아래 레이아웃 수치(inch)가 전부 어긋난다.
SRC = os.path.join(HERE, "_theme_base.pptx")
BASE = os.path.join(HERE, "_theme_src.pptx")
TMP = os.path.join(HERE, "_build.pptx")
OUT = os.path.join(os.path.dirname(HERE), "ChainRiposte.pptx")

# ─────────────────────────────────────────── 디자인 토큰 (원본 실측값)
INK    = RGBColor(0x00, 0x00, 0x00)
BODY   = RGBColor(0x21, 0x21, 0x21)
MUTED  = RGBColor(0x5C, 0x5C, 0x5C)
DARK   = RGBColor(0x1D, 0x1D, 0x1D)
BLUE   = RGBColor(0x00, 0x73, 0xFF)
BLUE_D = RGBColor(0x00, 0x4C, 0xB3)
BLUE_L = RGBColor(0x7A, 0xB8, 0xFF)
BLUE_T = RGBColor(0xEC, 0xF3, 0xFF)
CARD   = RGBColor(0xEB, 0xEB, 0xEB)
SOFT   = RGBColor(0xF5, 0xF5, 0xF5)
LINE   = RGBColor(0xDD, 0xDD, 0xDD)
RULE   = RGBColor(0x3F, 0x3F, 0x3F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x6A, 0x6A, 0x6A)
GREY_L = RGBColor(0x99, 0x99, 0x99)
RED    = RGBColor(0xD1, 0x36, 0x2B)
RED_T  = RGBColor(0xFD, 0xF0, 0xEF)

F_B, F_L, F_SB, F_EB, F_C = ("Pretendard Bold", "Pretendard Light",
                             "Pretendard SemiBold", "Pretendard ExtraBold", "Consolas")

# 레이아웃 (inch)
PANEL_X, PANEL_W = 0.78, 25.12
HDR_Y, HDR_H     = 0.74, 0.93
BODY_Y, BODY_H   = 1.74, 12.53
CX, CW           = 2.24, 22.19
TITLE_Y, TITLE_H = 2.06, 1.35
RULE1_Y, SUB_Y, RULE2_Y = 3.64, 4.02, 4.93
TOP_SUB, TOP_NOSUB = 5.40, 4.14
BOT_RULE_Y, BOTTOM = 13.62, 13.45

DATE, BRAND = "2026-08", "ChainRiposte"
TOKEN = re.compile(r"\*\*(.+?)\*\*|`(.+?)`")
MIN = 0.12  # 이보다 작은 도형은 만들지 않는다


class V:
    """세로 공간 배분기. 남은 공간을 넘겨 쓰면 즉시 죽는다."""

    def __init__(self, top, bottom, label=""):
        self.y, self.bottom, self.label = top, bottom, label

    def take(self, h, gap=0.0):
        assert h >= MIN, f"{self.label}: 높이 {h:.2f} 가 너무 작다"
        assert self.y + h <= self.bottom + 1e-6, (
            f"{self.label}: 세로 공간 초과 — {h:.2f} 를 넣으려는데 {self.bottom - self.y:.2f} 남음")
        y = self.y
        self.y = y + h + gap
        return y

    def rest(self):
        r = self.bottom - self.y
        assert r >= MIN, f"{self.label}: 남은 공간 없음 ({r:.2f})"
        return r

    def take_rest(self):
        r = self.rest()
        return self.take(r), r


# ─────────────────────────────────────────── 저수준 헬퍼
def _set_font(run, name, size, color, bold=False):
    f = run.font
    f.name, f.bold = name, bold
    f.size = Pt(size)
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def _effect(shape, xml):
    spPr = shape._element.spPr
    old = spPr.find(qn("a:effectLst"))
    if old is not None:
        spPr.remove(old)
    spPr.append(etree.fromstring(xml))


NO_SHADOW = '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
SHADOW = ('<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
          '<a:outerShdw blurRad="76200" dist="25400" dir="5400000" rotWithShape="0">'
          '<a:srgbClr val="000000"><a:alpha val="9000"/></a:srgbClr></a:outerShdw></a:effectLst>')


def rect(sl, x, y, w, h, fill=None, line=None, lw=1.0, radius=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=False, dash=False):
    assert w >= MIN and h >= MIN, f"도형 크기 비정상 w={w:.2f} h={h:.2f} (x={x:.2f} y={y:.2f})"
    s = sl.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = max(0.0, min(0.5, radius / min(w, h)))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
        if dash:
            s.line.dash_style = 4
    s.shadow.inherit = False
    _effect(s, SHADOW if shadow else NO_SHADOW)
    s.text_frame.word_wrap = True
    return s


def hline(sl, x, y, w, color=RULE, lw=1.25):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y),
                                Inches(x + w), Inches(y))
    c.line.color.rgb, c.line.width = color, Pt(lw)
    return c


def vline(sl, x, y, h, color=RULE, lw=1.25):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y),
                                Inches(x), Inches(y + h))
    c.line.color.rgb, c.line.width = color, Pt(lw)
    return c


def _fill_tf(tf, paras, size, font, color, align, spacing, em, space_after):
    if isinstance(paras, str):
        paras = [paras]
    first = True
    for item in paras:
        opt = dict(item) if isinstance(item, dict) else {}
        s = opt.pop("t", "") if isinstance(item, dict) else item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opt.get("align", align)
        p.line_spacing = opt.get("spacing", spacing)
        if opt.get("space_before"):
            p.space_before = Pt(opt["space_before"])
        sa = opt.get("space_after", space_after)
        if sa:
            p.space_after = Pt(sa)
        f_, sz = opt.get("font", font), opt.get("size", size)
        co, emc = opt.get("color", color), opt.get("em", em)
        if not s:
            r = p.add_run(); r.text = " "; _set_font(r, f_, sz, co); continue
        pos = 0
        for m in TOKEN.finditer(s):
            if m.start() > pos:
                r = p.add_run(); r.text = s[pos:m.start()]; _set_font(r, f_, sz, co)
            if m.group(1) is not None:
                r = p.add_run(); r.text = m.group(1); _set_font(r, F_B, sz, emc)
            else:
                r = p.add_run(); r.text = m.group(2); _set_font(r, F_C, sz * 0.92, BLUE_D)
            pos = m.end()
        if pos < len(s):
            r = p.add_run(); r.text = s[pos:]; _set_font(r, f_, sz, co)


def text(sl, x, y, w, h, paras, size=27, font=F_B, color=BODY, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.22, em=BLUE_D, space_after=0.0, box=None,
         fit=True):
    assert w >= MIN and h >= MIN, f"텍스트 상자 비정상 w={w:.2f} h={h:.2f} (y={y:.2f})"
    tb = box or sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if not fit:
        tb.name = "FREE " + tb.name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    _fill_tf(tf, paras, size, font, color, align, spacing, em, space_after)
    return tb


def label_in(shape, s, size, font=F_B, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, em=None, spacing=1.15, pad=0.18):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(pad)
    tf.margin_top = tf.margin_bottom = 0
    _fill_tf(tf, s, size, font, color, align, spacing, em or color, 0.0)
    return shape


# ─────────────────────────────────────────── 슬라이드 골격
def chrome(prs, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, PANEL_X, BODY_Y, PANEL_W, BODY_H, fill=WHITE, radius=0.14, shadow=True)
    rect(sl, PANEL_X, HDR_Y, PANEL_W, HDR_H, fill=DARK, radius=0.12)
    text(sl, PANEL_X + 0.71, HDR_Y + 0.20, 6.0, 0.55, BRAND, size=25.5, color=WHITE)
    text(sl, PANEL_X + PANEL_W - 7.0, HDR_Y + 0.26, 6.3, 0.45, DATE, size=20,
         color=WHITE, align=PP_ALIGN.RIGHT)
    hline(sl, 1.46, BOT_RULE_Y, 23.26, RULE, 1.1)
    text(sl, PANEL_X + PANEL_W - 6.5, BOT_RULE_Y + 0.13, 5.9, 0.45, str(page),
         size=18.5, color=DARK, align=PP_ALIGN.RIGHT)
    return sl


def head(sl, title, sub=None):
    text(sl, 4.21, TITLE_Y, 18.25, TITLE_H, title, size=67.5, color=INK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, fit=False)
    hline(sl, 1.44, RULE1_Y, 23.79, RULE, 1.25)
    if sub:
        text(sl, 2.6, SUB_Y, 21.5, 0.62, sub, size=23, font=F_L, color=BODY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, fit=False)
        hline(sl, 1.44, RULE2_Y, 23.79, RULE, 1.25)
        return V(TOP_SUB, BOTTOM, title)
    return V(TOP_NOSUB, BOTTOM, title)


# ─────────────────────────────────────────── 컴포넌트
def badge(sl, x, y, w, h, s, size=25, fill=BLUE, color=WHITE):
    return label_in(rect(sl, x, y, w, h, fill=fill, radius=0.10), s, size, F_B, color)


def row_card(sl, y, h, tag, title, desc, tag_w=3.2, title_w=6.2, fill=SOFT,
             size_title=31, size_desc=26.5):
    rect(sl, CX, y, CW, h, fill=fill, radius=0.13, line=LINE, lw=0.75)
    badge(sl, CX + 0.55, y + (h - 0.72) / 2, tag_w, 0.72, tag, size=24)
    tx = CX + 0.55 + tag_w + 0.55
    text(sl, tx, y + 0.2, title_w, h - 0.4, title, size=size_title, color=INK,
         anchor=MSO_ANCHOR.MIDDLE)
    dx = tx + title_w + 0.5
    text(sl, dx, y + 0.2, CX + CW - 0.6 - dx, h - 0.4, desc, size=size_desc,
         color=BODY, anchor=MSO_ANCHOR.MIDDLE)


def col_card(sl, x, y, w, h, kicker, title, lines, kicker_fill=BLUE,
             size_title=34, size_body=26, fill=SOFT, pad=0.45):
    rect(sl, x, y, w, h, fill=fill, radius=0.13, line=LINE, lw=0.75)
    cy = y + pad
    if kicker:
        badge(sl, x + pad, cy, min(4.2, w - 2 * pad), 0.64, kicker, size=22, fill=kicker_fill)
        cy += 0.64 + 0.18
    text(sl, x + pad, cy, w - 2 * pad, 0.72, title, size=size_title, color=INK)
    cy += 0.72 + 0.10
    bh = (y + h - pad) - cy
    text(sl, x + pad, cy, w - 2 * pad, bh, lines, size=size_body, color=BODY,
         spacing=1.30, space_after=7)


def callout(sl, y, h, quote, note=None, fill=BLUE_T, bar=BLUE, size_q=33, size_n=26,
            qcolor=BLUE_D):
    rect(sl, CX, y, CW, h, fill=fill, radius=0.13)
    rect(sl, CX, y, 0.17, h, fill=bar, shape=MSO_SHAPE.RECTANGLE)
    paras = [{"t": quote, "size": size_q, "color": qcolor}]
    if note:
        paras.append({"t": note, "size": size_n, "color": BODY, "space_before": 10})
    text(sl, CX + 0.85, y + 0.22, CW - 1.7, h - 0.44, paras,
         anchor=MSO_ANCHOR.MIDDLE, em=qcolor)


def table(sl, y, h, cols, header, rows, hdr_h=0.80, gap=0.13, highlight=None,
          size_h=24, size_r=26.5):
    """h 안에 머리행 + 본문행을 균등 배분한다."""
    n = len(rows)
    row_h = (h - hdr_h - gap * n) / n
    assert row_h >= 0.5, f"표 행 높이가 너무 작다 ({row_h:.2f})"
    tot = sum(cols)
    xs, acc = [], CX
    for c in cols:
        xs.append((acc, CW * c / tot))
        acc += CW * c / tot
    rect(sl, CX, y, CW, hdr_h, fill=DARK, radius=0.09)
    for (x, w), t in zip(xs, header):
        text(sl, x + 0.55, y, w - 1.0, hdr_h, t, size=size_h, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    yy = y + hdr_h + gap
    for i, row in enumerate(rows):
        hot = (highlight is not None and i == highlight)
        rect(sl, CX, yy, CW, row_h,
             fill=(BLUE_T if hot else (SOFT if i % 2 == 0 else CARD)),
             radius=0.09, line=(BLUE if hot else None), lw=1.5)
        for j, ((x, w), t) in enumerate(zip(xs, row)):
            text(sl, x + 0.55, yy + 0.06, w - 1.0, row_h - 0.12, t, size=size_r,
                 color=(BLUE_D if hot and j == 0 else (INK if j == 0 else BODY)),
                 anchor=MSO_ANCHOR.MIDDLE, em=(BLUE_D if hot else BLUE_D))
        yy += row_h + gap


def chips(sl, y, h, items, gap=0.30, fill=SOFT, color=INK, size=25, arrows=False,
          accent_last=False):
    n = len(items)
    aw = 0.62 if arrows else 0.0
    w = (CW - gap * (n - 1) - aw * (n - 1)) / n
    x = CX
    for i, it in enumerate(items):
        hot = accent_last and i == n - 1
        b = rect(sl, x, y, w, h, fill=(BLUE if hot else fill),
                 radius=0.11, line=(None if hot else LINE), lw=0.75)
        label_in(b, it, size, F_B, WHITE if hot else color)
        x += w
        if arrows and i < n - 1:
            text(sl, x, y, aw, h, "→", size=30, color=BLUE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
            x += aw
        x += gap


def media(sl, x, y, w, h, caption, sub=None):
    s = rect(sl, x, y, w, h, fill=SOFT, radius=0.12,
             line=RGBColor(0xB8, 0xB8, 0xB8), lw=1.6, dash=True)
    paras = [{"t": caption, "size": 26, "color": GREY}]
    if sub and h > 0.95:
        paras.append({"t": sub, "size": 21, "color": GREY_L, "font": F_L, "space_before": 7})
    label_in(s, paras, 26, F_B, GREY)
    return s


def stat(sl, x, y, w, h, big, cap, color=BLUE):
    rect(sl, x, y, w, h, fill=SOFT, radius=0.13, line=LINE, lw=0.75)
    text(sl, x + 0.3, y + 0.30, w - 0.6, h * 0.50, big, size=56, font=F_EB, color=color,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(sl, x + 0.3, y + h * 0.55, w - 0.6, h * 0.38, cap, size=23, color=BODY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def cols3(v, h, gap=0.44):
    """3열 카드 자리 (x, y, w, h) 목록 + 배정된 y."""
    w = (CW - gap * 2) / 3
    y = v.take(h, gap)
    return [(CX + (w + gap) * i, y, w, h) for i in range(3)], y


def cols2(v, h, gap=0.55):
    w = (CW - gap) / 2
    y = v.take(h, gap)
    return [(CX, y, w, h), (CX + w + gap, y, w, h)], y


def codebox(sl, x, y, w, h, lines, size=23, fill=DARK):
    """코드/의사코드 상자. 넘침 자동 축소 대상에서 빼려면 호출부에서 높이를 넉넉히 준다."""
    rect(sl, x, y, w, h, fill=fill, radius=0.10)
    paras = []
    for ln in lines:
        item = dict(ln) if isinstance(ln, dict) else {"t": ln}
        item.setdefault("font", F_C)
        item.setdefault("color", RGBColor(0xDC, 0xE6, 0xF5))
        item.setdefault("size", size)
        paras.append(item)
    text(sl, x + 0.55, y + 0.26, w - 1.1, h - 0.52, paras, size=size, font=F_C,
         color=RGBColor(0xDC, 0xE6, 0xF5), spacing=1.34, em=BLUE_L)


class Column:
    """본문 폭(CW)을 잠시 좁힌다 — 오른쪽 미디어 레일을 깔 때 왼쪽 칼럼용."""

    def __init__(self, width):
        self.width = width

    def __enter__(self):
        global CW
        self._prev = CW
        CW = self.width
        return self

    def __exit__(self, *exc):
        global CW
        CW = self._prev
        return False


def rail(sl, v, caption, sub=None, width=8.60, gap=0.55):
    """
    오른쪽에 <b>본문 높이를 통째로 쓰는</b> 스크린샷/영상 자리를 깔고,
    남은 왼쪽 폭으로 CW를 좁힌 컨텍스트를 돌려준다.

    작은 자리를 아래에 붙이면 도판이 장식이 된다 — 발표에서 실제로 보여 줄 것은
    화면이므로 자리부터 크게 잡고 글을 그 옆에 맞춘다.
    """
    top = v.y
    media(sl, CX + CW - width, top, width, BOTTOM - top, caption, sub)
    return Column(CW - width - gap)


def lean_card(sl, x, y, w, h, title, lines, size_title=30, size_body=24, fill=SOFT):
    """제목 한 줄 + 본문. col_card보다 여백이 얇아 좁은 칼럼에서 본문 자리가 남는다."""
    pad = 0.40
    rect(sl, x, y, w, h, fill=fill, radius=0.13, line=LINE, lw=0.75)
    text(sl, x + pad, y + pad, w - 2 * pad, 0.56, title, size=size_title, color=INK)
    text(sl, x + pad, y + pad + 0.62, w - 2 * pad, h - 2 * pad - 0.62, lines,
         size=size_body, color=BODY, spacing=1.28, space_after=7)


def stack(v, h, gap=0.30):
    """왼쪽 칼럼에 카드를 한 장 쌓을 자리 (x, y, w, h)."""
    return CX, v.take(h, gap), CW, h


def index_grid(sl, x, y, cols, rows, hot=None, cell=0.62, gap=0.07, size=20):
    """행 우선 평탄화(index = row × width + x)를 눈으로 보여 주는 번호 격자."""
    step = cell + gap
    for r in range(rows):
        for c in range(cols):
            i = r * cols + c
            on = (hot is not None and i == hot)
            s = rect(sl, x + c * step, y + r * step, cell, cell,
                     fill=(BLUE if on else CARD), line=(None if on else LINE), lw=0.7,
                     radius=0.06, shape=MSO_SHAPE.RECTANGLE)
            label_in(s, str(i), size, F_B, WHITE if on else MUTED, pad=0.0)
    return cols * step - gap, rows * step - gap


# 미니 보드 도식 — 알고리즘을 말로만 적으면 안 읽힌다
CELL_TILE, CELL_WALL, CELL_HOLE, CELL_PATH, CELL_SKIP = "o", "W", ".", "*", " "


def mini_board(sl, x, y, rows, cell=0.44, gap=0.05, caption=None, cap_size=21):
    """문자 격자를 도형으로 그린다. o=타일 · W=벽 · .=빈 칸 · *=밀려날 타일 · 공백=없음."""
    step = cell + gap
    for r, row in enumerate(rows):
        for c, ch in enumerate(row):
            if ch == CELL_SKIP:
                continue
            cx, cy = x + c * step, y + r * step
            if ch == CELL_WALL:
                s = rect(sl, cx, cy, cell, cell, fill=DARK, radius=0.05,
                         shape=MSO_SHAPE.RECTANGLE)
                label_in(s, "벽", 14, F_B, WHITE, pad=0.0)
            elif ch == CELL_HOLE:
                rect(sl, cx, cy, cell, cell, fill=RED_T, line=RED, lw=1.6, radius=0.05,
                     shape=MSO_SHAPE.RECTANGLE, dash=True)
            elif ch == CELL_PATH:
                rect(sl, cx, cy, cell, cell, fill=BLUE_L, line=BLUE, lw=1.2, radius=0.05,
                     shape=MSO_SHAPE.RECTANGLE)
            else:
                rect(sl, cx, cy, cell, cell, fill=CARD, line=LINE, lw=0.7, radius=0.05,
                     shape=MSO_SHAPE.RECTANGLE)

    w = len(rows[0]) * step - gap
    h = len(rows) * step - gap
    if caption:
        text(sl, x - 0.3, y + h + 0.14, w + 0.6, 0.42, caption, size=cap_size, font=F_L,
             color=MUTED, align=PP_ALIGN.CENTER)
    return w, h


# ═══════════════════════════════════════════ 슬라이드
def s01_cover(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0.76, 0.74, 1.11, 13.53, fill=BLUE, radius=0.12)
    tb = sl.shapes.add_textbox(Inches(-4.39), Inches(6.6), Inches(5.25), Inches(0.9))
    text(sl, 0, 0, 1, 1, "Portfolio", size=44, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, box=tb)
    tb.rotation = 270
    tb2 = sl.shapes.add_textbox(Inches(-2.10), Inches(11.9), Inches(3.4), Inches(0.5))
    text(sl, 0, 0, 1, 1, DATE, size=20, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, box=tb2)
    tb2.rotation = 270

    rect(sl, 1.96, 0.74, 23.93, 13.53, fill=WHITE, radius=0.14, shadow=True)
    hline(sl, 2.79, 1.39, 20.29, RULE, 1.1)
    badge(sl, 23.20, 1.06, 1.90, 0.66, "1", size=20, fill=DARK)

    # 표지 문구는 사용자가 직접 고친 것을 그대로 따른다 (이름 40pt, 부제 없음).
    text(sl, 4.4, 4.18, 19.0, 2.30, "ChainRiposte", size=110, font=F_L, color=INK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, fit=False)
    text(sl, 4.4, 6.08, 19.0, 2.30, "양평화", size=40, font=F_EB,
         color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, fit=False)
    hline(sl, 12.51, 8.75, 2.87, BLUE, 3.0)
    chips(sl, 9.45, 1.05, ["Unity 6", "C#", "1인 개발", "3주 (2026.07.14 ~ 08.03)"], size=27)

    hline(sl, 2.78, 11.55, 22.18, RULE, 1.1)
    text(sl, 2.65, 12.05, 9.74, 0.5, "PROJECT", size=24, font=F_SB, color=BLUE,
         align=PP_ALIGN.CENTER)
    text(sl, 2.67, 12.67, 9.72, 0.9, "Unity 퍼즐 + 패링 전투 · 매치3 → 보스전 반복",
         size=22, font=F_L, color=BODY, align=PP_ALIGN.CENTER)
    vline(sl, 13.33, 12.00, 1.35, LINE, 1.2)
    text(sl, 14.26, 12.05, 9.74, 0.5, "SCALE", size=24, font=F_SB, color=BLUE,
         align=PP_ALIGN.CENTER)
    text(sl, 14.28, 12.67, 9.72, 0.9,
         "170파일 / 25,546줄 · 157커밋 · EditMode 테스트 167개 전부 통과",
         size=22, font=F_L, color=BODY, align=PP_ALIGN.CENTER)


def s02_game(prs):
    sl = chrome(prs, 2)
    v = head(sl, "두 장르를 한 자원으로 묶었다", "퍼즐은 준비 구간, 보스전은 실력 구간.")

    with rail(sl, v, "스크린샷 — 퍼즐 화면 + 보스전 화면",
              "세로 화면 2장을 위아래 또는 나란히", width=8.80):
        y = v.take(0.42, 0.10)
        text(sl, CX, y, CW, 0.42, "한 판의 흐름", size=25, color=BLUE)
        y = v.take(0.92, 0.34)
        chips(sl, y, 0.92, ["월드맵", "퍼즐", "준비", "보스전", "다음 고리"],
              arrows=True, size=22, accent_last=True)

        lean_card(sl, *stack(v, 3.20), "두 장르가 이어지는 지점",
                  ["퍼즐에서 캔 소울로 레벨업하고, **퍼즐에서 맞은 HP가 그대로 보스전으로 넘어간다.**",
                   "죽으면 **사슬만 끊기고 성장은 남는다** — 빌드를 잃지 않으므로 다시 도전할 이유가 생긴다."],
                  size_title=33, size_body=25)

        y, h = v.take_rest()
        callout(sl, y, h,
                "난이도는 **실행(패링)**에서 오고, 성장은 **모서리만 깎는다.**",
                "성장이 난이도를 지우면 퍼즐이 무의미해지고, 아무것도 안 하면 성장이 무의미해진다.",
                size_q=30, size_n=24)


def s03_data(prs):
    sl = chrome(prs, 3)
    v = head(sl, "판 하나 = 에셋 하나",
             "밸런스 수치는 코드에 없다. 전부 ScriptableObject 에셋이고, 인스펙터에서 고친다.")

    with rail(sl, v, "스크린샷 — StageDataSO 인스펙터 전체",
              "접힌 곳 없이 펼쳐서 · Project 창의 Data 폴더도 함께", width=9.20):
        dh = 2.05
        y = v.take(dh, 0.30)
        rect(sl, CX, y, CW, dh, fill=SOFT, radius=0.13, line=LINE, lw=0.75)
        bw, bh, ar = 3.45, 0.95, 0.62
        bx = CX + 0.42
        b = rect(sl, bx, y + 0.28, bw, bh, fill=WHITE, radius=0.10, line=BLUE, lw=1.8)
        label_in(b, "StageDataSO", 22, F_C, BLUE_D, pad=0.10)
        text(sl, bx + bw, y + 0.28, ar, bh, "→", size=28, color=BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        bx2 = bx + bw + ar
        b = rect(sl, bx2, y + 0.28, bw, bh, fill=BLUE, radius=0.10)
        label_in(b, "ToConfig()", 22, F_C, WHITE, pad=0.10)
        text(sl, bx2 + bw, y + 0.28, ar, bh, "→", size=28, color=BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        b = rect(sl, bx2 + bw + ar, y + 0.28, bw, bh, fill=WHITE, radius=0.10, line=BLUE, lw=1.8)
        label_in(b, "퍼즐 · 전투 엔진", 22, F_B, BLUE_D, pad=0.10)
        text(sl, CX + 0.42, y + 1.36, CW - 0.84, 0.50,
             "인스펙터가 다루기 좋은 값을 **엔진이 쓰는 형태로 번역하는 창구가 하나**다.",
             size=23, color=BODY)

        lean_card(sl, *stack(v, 2.85), "번역은 한 곳에서만",
                  ["보드는 `\"OOWOO\"` 문자열로 그린다 — git diff에 모양이 그대로 보인다.",
                   "마스크와 벽 좌표로 바꾸면서 **위아래를 뒤집는 곳이 딱 한 군데**다. 두 곳에서 뒤집으면 제자리다."],
                  size_title=31, size_body=24)

        y, h = v.take_rest()
        lean_card(sl, CX, y, CW, h, "늘리는 비용이 에셋 1개",
                  ["스테이지 · 캐릭터 · 보스 · 타일 — 전부 **에셋 하나 추가**로 붙는다.",
                   "SO **10종** · 데이터 에셋 **20여 개** · 코드에 상수로 박힌 밸런스 **0개**."],
                  size_title=31, size_body=24)


def s04_puzzle_flow(prs):
    sl = chrome(prs, 4)
    v = head(sl, "퍼즐 엔진 — 한 번의 스왑",
             "핵심 구현 ① — 엔진이 결과를 <b>한 번에 끝까지</b> 계산하고, 화면은 그 기록을 재생만 한다.")

    with rail(sl, v, "▶ 영상 5초 — 매치 → 연쇄 → 낙하",
              "한 번의 스왑이 끝까지 이어지는 장면", width=8.20):
        y = v.take(0.40, 0.10)
        text(sl, CX, y, CW, 0.40, "`TrySwap(a, b)` 안에서 도는 순서", size=24, color=BLUE)
        y = v.take(0.92, 0.32)
        chips(sl, y, 0.92, ["교환", "매치 찾기", "파괴", "중력 · 리필", "연쇄"],
              arrows=True, size=21, accent_last=True)

        lean_card(sl, *stack(v, 2.75), "매치 찾기 — 런을 모아 합친다",
                  ["가로 · 세로를 훑어 **3칸 이상 이어진 런**을 모으고, 겹치면 합친다(ㄱ자 · T자).",
                   "파괴 목록은 `HashSet`이라 **교차점이 두 번 세어지지 않는다.**"],
                  size_title=31, size_body=24)

        y, h = v.take_rest()
        lean_card(sl, CX, y, CW, h, "모델 / 뷰 — 완결된 기록을 넘긴다",
                  ["매치가 없으면 **그 자리에서 되돌리고 턴도 안 센다** — 그래서 수가 없으면 판이 멈춘다(다음 장).",
                   "`SwapResult` 하나에 **연쇄 · 낙하 · 기믹**이 전부 담긴다. 연출 중에도 **모델은 이미 최종 상태**라 건너뛰어도 결과가 같다."],
                  size_title=31, size_body=24)


def s05_deadlock(prs):
    sl = chrome(prs, 5)
    v = head(sl, "둘 수 있는 수를 전부 센다",
             "핵심 구현 ② — 매치 없는 스왑은 턴을 안 먹으므로, 수가 0이면 판이 영영 멈춘다.")

    boxes, y = cols2(v, 4.90, 0.35)
    col_card(sl, *boxes[0], None, "전수 조사 — 스왑 → 판정 → 되돌리기",
             ["후보는 **인접한 두 칸**뿐이다. 각 칸에서 **오른쪽 · 위 두 방향만** 보면 모든 쌍을 정확히 한 번 센다.",
              "실제로 바꿔 보고 판정한 뒤 **원래대로 되돌린다** — 탐색이 보드를 남기지 않는다.",
              "판정 기준을 엔진의 실제 스왑과 **같은 함수 하나**로 공유한다. 따로 적으면 «둘 수 있다고 했는데 못 두는» 유령 수가 생긴다."],
             size_title=33, size_body=24)
    col_card(sl, *boxes[1], None, "핵심 최적화 — 전체를 다시 훑지 않는다",
             ["**전제**: 정착이 끝난 보드에는 매치가 하나도 없다.",
              "그러면 새 매치는 반드시 **바뀐 두 칸 중 하나를 지난다** → 보드 전체 스캔 대신 그 칸에서 **네 방향으로 같은 종류를 세는** 국소 검사면 충분하다.",
              "**O(W²H²) → O(W·H·(W+H)).** 9×7 보드에서 약 **8,000 → 2,000 연산**이고, 매 스왑 끝마다 도는 검사다."],
             size_title=33, size_body=24)

    y, h = v.take_rest()
    rect(sl, CX, y, CW, h, fill=SOFT, radius=0.13, line=LINE, lw=0.75)
    text(sl, CX + 0.55, y + 0.24, CW - 1.1, 0.52,
         "수가 0이면 — 섞는다 (턴은 소모하지 않는다)", size=29, color=INK)
    text(sl, CX + 0.55, y + 0.90, CW - 1.1, h - 1.20,
         ["**움직일 수 있는 타일만** 모아(벽 · 보스 · 부패 · 사슬은 제자리) **Fisher–Yates**로 O(n) 균등 셔플한다. 난수를 주입받으므로 같은 씨앗이면 그대로 재현된다.",
          "섞자마자 매치가 터져 있으면 **공짜 콤보**가 되므로 될 때까지 **다시 뽑는다**(거부 표본추출, 최대 32회). 값이 아니라 **`Tile` 객체째** 옮기므로 폭탄 카운트가 따라간다."],
         size=24, color=BODY, spacing=1.26, space_after=6)


def s06_gravity(prs):
    sl = chrome(prs, 6)
    v = head(sl, "빈칸을 남기지 않는 중력",
             "핵심 구현 ③ — 벽과 구멍이 뚫린 비정형 보드에서도 성립해야 하는 규칙.")

    y = v.take(1.22, 0.28)
    callout(sl, y, 1.22,
            "불변식 — 정착이 끝나면 **구멍(X) 이외의 모든 칸이 채워져 있다.**", size_q=31)

    ch = 4.40
    y = v.take(ch, 0.28)
    rect(sl, CX, y, CW, ch, fill=SOFT, radius=0.13, line=LINE, lw=0.75)

    bx, by = CX + 0.62, y + 0.42
    before = ["*ooooo", "*ooooo", "*WWWoo", "o*.ooo", "oooooo"]
    after = [".ooooo", "oooooo", "oWWWoo", "oooooo", "oooooo"]
    bw, bh = mini_board(sl, bx, by, before, caption="① 갇힌 칸 + 길")
    text(sl, bx + bw + 0.12, by, 0.86, bh, "→", size=34, color=BLUE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    mini_board(sl, bx + bw + 0.98, by, after, caption="② 한 칸씩 밀었다")

    tx = bx + bw * 2 + 1.60
    text(sl, tx, y + 0.40, CX + CW - 0.7 - tx, ch - 0.8,
         ["**① 직선 낙하** — 열 안에서 아래로 압축한다. 구멍은 통과하고 벽 위에는 쌓인다.",
          "**② 대각선 슬라이드** — 벽 그늘에 든 칸은 대각선 위에서 끌어내리고, 벽에 얹혀 못 내려가는 타일은 대각선 아래로 흘려보낸다.",
          "**③ 리필** — 각 열 꼭대기의 **연속된** 빈 칸만 새 타일로 채운다(중간 빈 칸은 다음 웨이브의 낙하 몫이다).",
          "**④ 갇힌 칸 끌어오기** — ①~③이 전부 멈췄는데 빈 칸이 남으면, 그 칸에서 **리필이 닿는 자리까지 8방향 너비 우선 탐색**으로 최단 경로를 찾아 길 위의 타일을 한 칸씩 민다. **구멍이 반대로 꼭대기까지 걸어 올라간다.**"],
         size=24, color=BODY, spacing=1.28, space_after=8)

    y, h = v.take_rest()
    row_card(sl, y, h, "실측", "99.08%",
             "갇힌 칸의 **99.08%가 끌어오기로** 메워진다 · 제자리 생성은 **0.92%**(벽 · 보드 끝에 완전히 둘러싸여 길조차 없는 칸)  ·  경로 탐색은 `Queue` + 「어디서 왔는가」 사전으로 되짚어 만든다",
             tag_w=2.2, title_w=4.0, size_desc=24)


def s07_board_data(prs):
    sl = chrome(prs, 7)
    v = head(sl, "2차원 격자를 저장하는 법",
             "핵심 구현 ④ — 유니티는 `bool[,]`를 직렬화하지 못한다. 그래서 <b>눕혀야</b> 한다.")

    ch = 4.35
    y = v.take(ch, 0.30)
    rect(sl, CX, y, CW, ch, fill=SOFT, radius=0.13, line=LINE, lw=0.75)

    gx, gy = CX + 0.70, y + 0.55
    gw, gh = index_grid(sl, gx, gy, 5, 3, hot=7)
    text(sl, gx - 1.0, gy + gh + 0.18, gw + 2.0, 0.50,
         "width = 5 · `boardRows[1][2]` → **7**", size=21, font=F_L, color=MUTED,
         align=PP_ALIGN.CENTER)

    tx = gx + gw + 1.10
    text(sl, tx, y + 0.42, CX + CW - 0.7 - tx, ch - 0.84,
         ["**문제** — 유니티 직렬화기는 다차원 배열도 중첩 배열도 저장하지 못한다. 인스펙터에 격자를 띄우려면 **데이터가 먼저 저장 가능**해야 한다.",
          "**표준 해법(평탄화)** — 2차원을 1차원으로 눕히고 번호로 좌표를 되찾는다. **행 우선**으로 `index = row × width + x`, 역산은 `row = i / width` · `x = i % width` — 그래서 **가로 크기를 반드시 같이 저장**해야 한다.",
          "**이 프로젝트** — 행마다 문자열 하나(`string[] boardRows`). `boardRows[row][x]`는 **같은 행 우선 배치**이고, 다만 그 곱셈을 **문자열 컨테이너가 대신 들고 있다.**"],
         size=24, color=BODY, spacing=1.28, space_after=8)

    boxes, y = cols2(v, 3.40, 0.40)
    col_card(sl, *boxes[0], None, "왜 리스트가 아니라 문자열인가",
             ["칸이 **3상태**(O · X · W)라 `List<bool>`로는 못 담는다.",
              "git diff에 **`\"OOWOO\"`로 사람이 읽힌다.** 원소 63개 리스트는 안 읽힌다.",
              "만들어 둔 **스테이지 8개를 옮기지 않아도** 된다."],
             size_title=33, size_body=24)
    col_card(sl, *boxes[1], None, "번호는 이럴 때 쓴다",
             ["드래그로 칠할 때 **마지막에 칠한 칸**을 번호 하나로 비교한다.",
              "런타임은 반대로 편다 — 엔진은 순수 C#이라 **진짜 2차원 `bool[x, y]`**다."],
             size_title=33, size_body=24)


def s08_board_gui(prs):
    sl = chrome(prs, 8)
    v = head(sl, "격자를 그리고 칠하기",
             "핵심 구현 ⑤ — 인스펙터의 격자는 위젯이 아니라 <b>직접 그린 사각형</b>이다.")

    with rail(sl, v, "스크린샷 — 인스펙터 격자 + Board Editor 창",
              "같은 스테이지를 두 곳에서 편집하는 화면", width=9.20):
        lean_card(sl, *stack(v, 2.55), "그리기 — 즉시 모드 GUI",
                  ["`GetRect`로 자리 하나를 받고, 칸의 사각형을 **직접 계산해 채운다.**",
                   "버튼 63개를 만드는 게 아니라 **이벤트마다 다시 그린다.**"],
                  size_title=30, size_body=24)
        lean_card(sl, *stack(v, 2.55), "칠하기 — 좌표를 번호로",
                  ["마우스가 어느 칸인지는 같은 식의 **역산**이다.",
                   "**같은 칸은 한 드래그에서 한 번만** 처리한다 — 안 막으면 Undo 기록이 수십 개 쌓인다."],
                  size_title=30, size_body=24)

        y, h = v.take_rest()
        callout(sl, y, h, "그리는 코드는 한 벌, 편집 상태는 각자",
                "인스펙터와 전용 창이 **같은 그리기 코드**를 부르고, 「어느 칸을 칠하는 중인가」만 **창마다 하나씩** 들고 넘긴다 — `static`이면 서로 덮어써 드래그가 끊긴다.",
                size_q=28, size_n=23)


def s09_combat_chart(prs):
    sl = chrome(prs, 9)
    v = head(sl, "리듬 패링 전투 — 채보와 판정",
             "핵심 구현 ⑥ — 보스 공격은 애니메이션이 아니라 <b>악보</b>다.")

    with rail(sl, v, "▶ 영상 5초 — 패링 판정 링",
              "+ 인살 컷씬 · 정지 이미지로는 설명이 안 된다", width=8.20):
        lean_card(sl, *stack(v, 2.50, 0.26), "데이터 — 노트 = 값 3개",
                  ["**타격 시점(박) · 예비동작 길이(박) · 피해.** 패턴 = BPM + 길이 + 노트 목록.",
                   "**연속기는 별도 타입이 아니라 촘촘히 찍은 노트**다."],
                  size_title=30, size_body=24)
        lean_card(sl, *stack(v, 2.55, 0.26), "판정 — 누른 순간 결판",
                  ["`[타격 − 윈도우, 타격 + 유예]` 안의 노트 중 **가장 임박한 하나**를 그 자리에서 지운다.",
                   "밖에서 누르면 **헛침** — 잠기고 **보스가 체간을 되찾는다.**"],
                  size_title=30, size_body=24)

        y, h = v.take_rest()
        lean_card(sl, CX, y, CW, h, "표시 — 거리 = 남은 시간",
                  ["반지름을 **`1 + 남은시간 × 접근속도`**로 정한다 → 회색 띠의 **두께가 곧 판정 폭**이다.",
                   "**흰 원의 안쪽 테두리가 노트의 위치**이고, 띠에 겹친 동안이 성공 구간이다."],
                  size_title=30, size_body=24)


def s10_deterministic(prs):
    sl = chrome(prs, 10)
    v = head(sl, "프레임에 기대지 않는 시간",
             "핵심 구현 ⑦ — 같은 입력이면 프레임레이트가 달라도 <b>같은 결과</b>가 나온다.")

    y = v.take(1.85, 0.26)
    rect(sl, CX, y, CW, 1.85, fill=RED_T, radius=0.13)
    rect(sl, CX, y, 0.17, 1.85, fill=RED, shape=MSO_SHAPE.RECTANGLE)
    text(sl, CX + 0.90, y + 0.20, CW - 1.8, 1.45,
         [{"t": "문제 — 매 프레임 `dt`만큼 통째로 굴리면", "size": 30, "color": RED},
          {"t": "타격 · 판정 만료 · 예비동작 시작이 **한 프레임 안에 뭉쳐** 순서가 뒤바뀌고, 처리도 최대 한 프레임 늦는다.",
           "size": 25, "color": BODY, "space_before": 8}],
         anchor=MSO_ANCHOR.MIDDLE, em=RED)

    y = v.take(2.62, 0.26)
    codebox(sl, CX, y, CW, 2.62,
            ["while (remaining > 0) {",
             "    step = min(remaining, 다음_사건까지, 플레이어_타이머);",
             "    Advance(step);        // 경계까지만 시간을 흘린다",
             "    remaining -= step;",
             "}"], size=21)

    boxes, y = cols2(v, 3.05, 0.40)
    col_card(sl, *boxes[0], None, "방법 — 사건 경계까지만",
             ["`dt`를 통째로 흘리지 않고 **다음 사건**(예비동작 시작 · 타격 · 유예 만료 · 패턴 종료) **직전까지만** 잘라 진행한다."],
             size_title=33, size_body=24)
    col_card(sl, *boxes[1], None, "효과 — 전투를 테스트한다",
             ["시간이 `Tick(초)`로만 흐르고 `UnityEngine`을 안 쓰므로, 패링 · 헛침 · 인살 · 페이즈 전환을 **엔진 없이** 검증한다."],
             size_title=33, size_body=24)


def s11_chart_editor(prs):
    sl = chrome(prs, 11)
    v = head(sl, "리듬(채보) 에디터",
             "핵심 구현 ⑧ — 보스 공격을 <b>그려서</b> 만들고, 플레이하지 않고 <b>들어 본다.</b>")

    with rail(sl, v, "스크린샷 — 채보 타임라인 + ▶ 미리듣기 무대",
              "가능하면 영상 5초 (재생 줄이 흐르고 원이 다가오는 장면)", width=9.20):
        lean_card(sl, *stack(v, 2.55), "타임라인 — 클릭해서 찍는다",
                  ["클릭 = 추가, 드래그 = 이동, 우클릭 = 삭제. 스냅 **1 · ½ · ¼박**.",
                   "예비동작은 **왼쪽으로 뻗는 막대** — 언제부터 보이는지가 읽힌다."],
                  size_title=30, size_body=24)
        lean_card(sl, *stack(v, 2.85), "▶ 미리듣기 — 플레이 없이 돌린다",
                  ["`EditorApplication.update`로 **에디터 시간**을 굴려 원 · 판정 띠 · 타격음을 재생한다.",
                   "판정 폭은 **`PlayerStatsConfigSO`의 실제 값**을 읽는다 — 제 숫자를 들면 채보가 게임에서 어긋난다."],
                  size_title=30, size_body=24)

        y, h = v.take_rest()
        callout(sl, y, h, "확인에 드는 시간: **2~3분 → 2초**",
                "없을 때는 찍기 → 플레이 → 퍼즐 클리어 → 보스전 도달까지 가야 한 번 들었다. 채보 튜닝은 **수십 번 반복하는 일**이다.",
                size_q=29, size_n=23)


def s12_loc(prs):
    sl = chrome(prs, 12)
    v = head(sl, "현지화 — 원천은 한 장",
             "원천을 한 곳으로 묶으면, 성능 최적화까지 정확해진다.")

    boxes, y = cols3(v, 4.75, 0.36)
    col_card(sl, *boxes[0], "구조", "CSV 한 장이 원천",
             ["구글 시트 → CSV(**139키 × 2언어**). 코드 · 씬에 문자열 하드코딩 **0**.",
              "정적 문구는 `LocalizedText` 바인더 + 키, 동적 문구(HP · 턴)는 `Loc.GetText(key, args)`."])
    col_card(sl, *boxes[1], "함정", "값이 아니라 키를 바꾼다",
             ["동적 문구에 바인더를 붙이면 언어 전환 시 **코드와 컴포넌트가 서로 덮어쓴다.**",
              "같은 칸을 여러 문구가 돌려 쓸 때는 값이 아니라 **키를 갈아 끼운다**(`SetKey`)."])
    col_card(sl, *boxes[2], "툴", "빠진 키를 찾아 준다",
             ["**CSV에 없는 키**를 쓰는 텍스트를 씬에서 훑어 잡는 메뉴.",
              "원천이 한 장뿐이라 **없는 글자는 게임에도 안 나온다** → 미리 구울 글자를 정확히 센다(**343자**)."])

    y, h = v.take_rest()
    gap, gw = 0.44, (CW - 0.88) / 3
    stat(sl, CX, y, gw, h, "139키 × 2언어", "원천 = CSV 한 장")
    stat(sl, CX + gw + gap, y, gw, h, "0", "코드 · 씬 하드코딩 문자열")
    stat(sl, CX + (gw + gap) * 2, y, gw, h, "343자", "시작 로딩에 미리 굽는 글자")


def s13_fuzz(prs):
    sl = chrome(prs, 13)
    v = head(sl, "보드 60,000개를 자동 검증했다",
             "엔진 규칙에 `UnityEngine`이 없어서, 유니티를 켜지 않고 검사할 수 있었다.")

    with rail(sl, v, "스크린샷 — 퍼즈 콘솔 출력",
              "+ Test Runner 167개 전부 초록", width=9.40):
        lean_card(sl, *stack(v, 4.90), "상황 → 방법 → 결과",
                  ["**상황** — 주석은 \"구멍 이외의 모든 칸이 항상 채워진다\"고 **주장**했고 **테스트 166개가 통과 중**이었다. 그런데 실제로는 벽 밑에 빈 칸이 남았다.",
                   "**방법** — 엔진 규칙이 **순수 C#**이라 콘솔 프로젝트에 소스를 그대로 물려 컴파일된다. 무작위 보드 **60,000개**를 돌려 반례를 모았다(테스트도 **`dotnet test` 0.3초** 완주).",
                   "**결과** — 이번 버그 **3개를 전부 퍼즈가 잡았다**(무한 루프 · 널 참조 · 빈 칸). 손으로 짠 테스트 166개는 **하나도 못 잡았다.**"],
                  size_title=31, size_body=24)

        y, h = v.take_rest()
        rect(sl, CX, y, CW, h, fill=DARK, radius=0.13)
        rect(sl, CX, y, 0.17, h, fill=BLUE, shape=MSO_SHAPE.RECTANGLE)
        text(sl, CX + 0.80, y + 0.20, CW - 1.6, h - 0.40,
             [{"t": "**불변식은 주장이지 보장이 아니다.**", "size": 30, "color": WHITE, "em": BLUE_L},
              {"t": "손으로 짠 예제는 «내가 생각한 판»만 검사한다. 생각하지 못한 판은 기계가 찾아야 했다.",
               "size": 24, "color": RGBColor(0xC9, 0xC9, 0xC9), "space_before": 7}],
             anchor=MSO_ANCHOR.MIDDLE)


def s14_trouble(prs):
    sl = chrome(prs, 14)
    v = head(sl, "트러블슈팅 대표 3건", "재발하는 버그는 주의가 아니라 코드로 막는다.")

    h = v.rest()
    boxes, y = cols3(v, h, 0.44)
    col_card(sl, *boxes[0], "① 확률의 단위", "보스 타일이 보드를 도배",
             ["난입 확률이 **리필 타일 하나하나에** 굴려지고 있었다. 한 웨이브에 10칸이 채워지면 주사위를 10번 던진다.",
              "**\"5%\"가 실제로는 `1 − 0.95¹⁰ ≈ 40%`였다.**",
              "→ 보드 위 동시 상한 + **웨이브 단위 누적 카운터**로 바꿨다.",
              "**확률은 \"무엇 당(per)\"인지 반드시 같이 적는다.**"],
             size_body=24)
    col_card(sl, *boxes[1], "② 모델은 맞고 뷰가 틀림", "타일이 옆 열을 가로질렀다",
             ["한 웨이브에서 두 번 움직인 타일(낙하 → 벽에서 대각 슬라이드)의 기록을 **`From → To` 하나로 합쳐** 뷰에 넘긴다 — 중간 좌표가 남으면 뷰가 칸을 잘못 추적하기 때문이다.",
              "그런데 뷰가 그 두 점을 **직선으로 이어** 타일이 벽을 뚫고 대각선으로 날았다.",
              "→ 뷰가 **꺾인 경로**로 재생하되, 속도가 꺾이는 지점에서 0이 되지 않게 **경로 전체 길이를 하나의 자**로 쓴다."],
             size_body=24)
    col_card(sl, *boxes[2], "③ 3번 재발", "게이지가 안 줄어든다",
             ["`Image.type`이 `Filled`가 아니면 `fillAmount`가 **아무 일도 안 한다.**",
              "UI 스프라이트를 꽂을 때 Unity가 9슬라이스를 감지해 `Sliced`로 되돌려 놓고 있었다.",
              "→ `Awake`에서 게이지 3종의 타입을 강제한다. **세 번 재발했다는 것은 「주의」로는 안 된다는 뜻이었다.**"],
             size_body=24)


def s15_opt(prs):
    sl = chrome(prs, 15)
    v = head(sl, "최적화 & 런타임 안정성",
             "최적화가 오히려 화면의 글자를 날린 사고 — 에디터에서는 멀쩡했다.")

    y = v.take(2.58, 0.28)
    rect(sl, CX, y, CW, 2.58, fill=RED_T, radius=0.13)
    rect(sl, CX, y, 0.17, 2.58, fill=RED, shape=MSO_SHAPE.RECTANGLE)
    text(sl, CX + 0.90, y + 0.20, CW - 1.8, 2.18,
         [{"t": "글자를 미리 굽게 했더니, 실기에서 일부가 안 보였다", "size": 30, "color": RED},
          {"t": "런타임 끊김을 없애려고 시작 로딩에서 폰트 글자를 **한 번에 구웠다.** 그런데 아틀라스가 **1024 한 장**이라 386자가 **두 장째로 넘어갔고**, 넘친 글자가 기기에서 안 그려졌다.",
           "size": 24, "color": BODY, "space_before": 8},
          {"t": "그전에는 **화면에 실제로 뜬 글자만** 조금씩 구워져 우연히 한 장에 들어 있었다 — **미리 굽는 행위 자체가 방아쇠**였다. → 아틀라스 **2048**로, 그리고 **두 장을 넘으면 에러 로그**를 남긴다.",
           "size": 24, "color": BODY, "space_before": 5}],
         anchor=MSO_ANCHOR.MIDDLE, em=RED)

    y = v.take(3.05, 0.28)
    table(sl, y, 3.05, [1.6, 2.0, 3.4], ["항목", "조치", "이유"],
          [["BGM 6곡", "Streaming / Vorbis", "원본 PCM 20MB × 6이 통째로 빌드에 들어가던 것 차단"],
           ["단발 효과음", "DecompressOnLoad", "스트리밍하면 **눌린 뒤에** 소리가 난다"],
           ["도트 스프라이트", "Point · 무압축 · 밉맵 끔", "이중선형 필터로 옆 타일보다 흐려 보이는 것 방지"]],
          hdr_h=0.74, gap=0.12, size_r=24.5)

    y, h = v.take_rest()
    rect(sl, CX, y, CW, h, fill=SOFT, radius=0.13, line=LINE, lw=0.75)
    text(sl, CX + 0.85, y + 0.18, CW - 1.7, h - 0.36,
         [{"t": "1인 개발에 특히 필요했던 안정성", "size": 28, "color": INK},
          {"t": "**널 세이프 슬롯**(클립 · 그림이 비어도 안 멈춘다)  ·  **도메인 리로드 OFF 대응** `ResetStatics()`  ·  세이브 **v1→v3** 마이그레이션  ·  정적 서비스 19종 · **`DontDestroyOnLoad` 2개**",
           "size": 24, "color": BODY, "space_before": 6}],
         anchor=MSO_ANCHOR.MIDDLE)


def s16_future(prs):
    sl = chrome(prs, 16)
    v = head(sl, "Future Work: 정직하게 남은 것", "시스템은 완성했고, 콘텐츠와 검증이 남았다.")

    y = v.take(4.95, 0.30)
    table(sl, y, 4.95, [2.2, 4.8], ["항목", "상태"],
          [["튜토리얼 2종", "설계 · 구현 완료 — 검증 · 영상 녹화 전이라 **현재 비활성화**"],
           ["효과음 3종", "슬롯 배선 완료, 클립 선정만 남음"],
           ["잡몹 밸런스 수치", "시스템 완료, 수치 입력 전 (공용값으로 동작 중)"],
           ["콘텐츠 볼륨", "월드 2개(6판). 시스템은 확장 대비 완료"]],
          hdr_h=0.78, gap=0.12)

    y = v.take(1.70, 0.26)
    rect(sl, CX, y, CW, 1.70, fill=SOFT, radius=0.13, line=LINE, lw=0.75)
    text(sl, CX + 0.85, y + 0.16, CW - 1.7, 1.38,
         [{"t": "남은 것이 적은 이유", "size": 29, "color": INK},
          {"t": "전부 **데이터를 채우는 일**이고 코드를 고칠 일이 아니다. 스테이지 · 캐릭터 · 언어 · 기믹 어느 쪽으로 늘려도 **에셋 1개 또는 상속 1개**로 붙는다.",
           "size": 25, "color": BODY, "space_before": 7}],
         anchor=MSO_ANCHOR.MIDDLE)

    y, h = v.take_rest()
    callout(sl, y, h,
            "튜토리얼이 꺼져 있다는 사실을 **먼저** 말한다 — 데모 중에 들통나는 것보다 낫다.",
            size_q=30)


def s17_retro(prs):
    sl = chrome(prs, 17)
    v = head(sl, "회고: 가장 크게 배운 것 3가지")

    def block(num, title, lines):
        h = 2.55
        y = v.take(h, 0.26)
        rect(sl, CX, y, CW, h, fill=SOFT, radius=0.13, line=LINE, lw=0.75)
        badge(sl, CX + 0.55, y + (h - 0.80) / 2, 0.80, 0.80, num, size=28)
        text(sl, CX + 1.65, y + 0.30, CW - 2.4, 0.58, title, size=31, color=INK)
        text(sl, CX + 1.65, y + 0.96, CW - 2.4, h - 1.26, lines, size=25,
             color=BODY, spacing=1.26, space_after=5)

    block("①", "같은 숫자를 두 곳에 적으면 반드시 어긋난다",
          ["패링 판정과 그것을 그리는 띠 · 채보 미리듣기와 실제 전투 · 튜토리얼의 밝은 칸과 누를 수 있는 칸 — 전부 같은 원인이었다. **한 곳에서 파생시키면 어긋날 수가 없다.**"])
    block("②", "불변식은 주장이지 보장이 아니다",
          ["\"모든 칸이 항상 채워진다\"는 주석과 통과 중인 테스트 166개가 오히려 안심시켰다. 규칙이 복잡해질수록 **사람이 짠 예제는 자기가 생각한 판만 검사한다** — 생각하지 못한 판은 기계(퍼즈)가 찾아야 했다."])
    block("③", "기획 의도를 코드가 배신하는 지점을 찾는 것이 진짜 일",
          ["잡몹 위협을 **턴**으로 세었더니 가만히 있는 것이 최적해가 됐다 → 실시간으로 변경.",
           "소울을 **클리어한 판에서 0**으로 했더니 빨리 깬 사람이 손해를 봤다 → 스테이지별 매장량으로."])

    y, h = v.take_rest()
    rect(sl, CX, y, CW, h, fill=DARK, radius=0.13)
    rect(sl, CX, y, 0.17, h, fill=BLUE, shape=MSO_SHAPE.RECTANGLE)
    text(sl, CX + 0.90, y, CW - 1.8, h,
         "돌아가게 만드는 것과 **「의도대로 돌아가게」** 만드는 것은 다른 일이었다.",
         size=35, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, em=BLUE_L)


# ═══════════════════════════════════════════ 실행
KEEP_FONTS = {"Pretendard Light", "Pretendard Bold", "Pretendard ExtraBold",
              "Pretendard SemiBold"}


def wipe_slides(prs):
    lst = prs.slides._sldIdLst
    for sld in list(lst):
        prs.part.drop_rel(sld.get(qn("r:id")))
        lst.remove(sld)


def strip_unused_fonts(prs):
    """AI 템플릿에서 딸려 온 일본어 폰트 3종(13MB)을 걷어낸다. 이 덱은 안 쓴다."""
    lst = prs.part._element.find(qn("p:embeddedFontLst"))
    if lst is None:
        return 0
    n = 0
    for ef in list(lst):
        fnt = ef.find(qn("p:font"))
        if fnt is not None and fnt.get("typeface") in KEEP_FONTS:
            continue
        for child in ef:
            rid = child.get(qn("r:id"))
            if rid:
                try:
                    prs.part.drop_rel(rid)
                except Exception:
                    pass
        lst.remove(ef)
        n += 1
    return n


SLIDES = [s01_cover, s02_game, s03_data,
          s04_puzzle_flow, s05_deadlock, s06_gravity,
          s07_board_data, s08_board_gui,
          s09_combat_chart, s10_deterministic, s11_chart_editor,
          s12_loc, s13_fuzz, s14_trouble, s15_opt,
          s16_future, s17_retro]


def fit_text(prs, floor=21.0, tol=1.02):
    """상자를 넘치는 글자를 줄인다. 'FREE ' 로 표시한 상자(주변이 여백인 큰 글씨)는 건드리지 않는다.
       PowerPoint 의 normAutofit 에 맡기지 않는 이유: 열어 보기 전에는 결과를 알 수 없다."""
    from validate import est_height
    shrunk = []
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame or sh.name.startswith("FREE "):
                continue
            tf = sh.text_frame
            if not tf.text.strip():
                continue
            w, h = Emu(sh.width).inches, Emu(sh.height).inches
            runs = [r for p in tf.paragraphs for r in p.runs if r.text]
            if not runs:
                continue
            start = max(r.font.size.pt for r in runs)
            scale = 1.0
            while est_height(tf, w) > h * tol and scale > 0.45:
                if max(r.font.size.pt for r in runs) <= floor:
                    break
                scale *= 0.95
                for r in runs:
                    r.font.size = Pt(max(floor, r.font.size.pt * 0.95))
                for p in tf.paragraphs:
                    if p.space_before:
                        p.space_before = Pt(p.space_before.pt * 0.95)
                    if p.space_after:
                        p.space_after = Pt(p.space_after.pt * 0.95)
            if scale < 1.0:
                shrunk.append((i, sh.name, start, max(r.font.size.pt for r in runs)))
    return shrunk


def build(out=TMP, slides=None, verbose=True):
    shutil.copyfile(SRC, BASE)
    prs = Presentation(BASE)
    wipe_slides(prs)
    strip_unused_fonts(prs)
    for fn in (slides or SLIDES):
        fn(prs)
    shrunk = fit_text(prs)
    if verbose and shrunk:
        print(f"자동 축소 {len(shrunk)}개 상자:")
        for i, n, a, b in shrunk:
            print(f"   슬라이드{i:2} {n:22} {a:.0f}pt → {b:.1f}pt")
    prs.save(out)
    return out


if __name__ == "__main__":
    out = build()
    if "--deploy" in sys.argv:
        shutil.copyfile(out, OUT)
        out = OUT
    print("built:", out, f"({len(SLIDES)}장)")
